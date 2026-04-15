# tests/test_powerpoint_output.py
import pytest
from pathlib import Path
from unittest.mock import patch


@pytest.fixture(autouse=True)
def isolated_projects_dir(tmp_path, monkeypatch):
    monkeypatch.setenv("PROJECTS_DIR", str(tmp_path))
    from api.config import get_settings
    get_settings.cache_clear()
    yield tmp_path
    get_settings.cache_clear()


@pytest.fixture
def slug(isolated_projects_dir):
    slug = "pptx-output-test"
    outputs_dir = isolated_projects_dir / slug / "outputs"
    outputs_dir.mkdir(parents=True)
    return slug


@pytest.fixture
def sample_slides():
    return [
        {
            "title": "Case for Change",
            "content": ["Manual processes slow delivery", "Data silos create errors"],
            "notes": "Emphasise cost of current state.",
        },
        {
            "title": "Our Value Propositions",
            "content": "Three high-value propositions identified.",
            "notes": "Refer to VP register.",
        },
    ]


@pytest.fixture
def sample_metadata():
    return {
        "org_name": "Acme Logistics Ltd",
        "financial_year": "FY2026",
        "sponsor": "Jane Smith, CEO",
        "date": "2026-04-15",
    }


def test_powerpoint_output_writes_file(slug, sample_slides, sample_metadata):
    """PowerPointOutputTool creates a .pptx file at the returned path."""
    from agents.tools.powerpoint_output import PowerPointOutputTool
    with patch("agents.tools.powerpoint_output.insert_agent_output_sync"):
        tool = PowerPointOutputTool(slug=slug)
        result = tool._run(
            slides=sample_slides,
            metadata=sample_metadata,
            filename="executive_presentation.pptx",
            agent_name="business_plan_generator",
        )
    assert Path(result).exists()


def test_powerpoint_output_returns_absolute_path(slug, sample_slides, sample_metadata):
    """Return value is an absolute path ending in .pptx."""
    from agents.tools.powerpoint_output import PowerPointOutputTool
    with patch("agents.tools.powerpoint_output.insert_agent_output_sync"):
        tool = PowerPointOutputTool(slug=slug)
        result = tool._run(
            slides=sample_slides,
            metadata=sample_metadata,
            filename="executive_presentation.pptx",
            agent_name="business_plan_generator",
        )
    assert Path(result).is_absolute()
    assert result.endswith(".pptx")


def test_powerpoint_output_contains_slide_titles(slug, sample_slides, sample_metadata):
    """Generated .pptx contains each slide title."""
    from agents.tools.powerpoint_output import PowerPointOutputTool
    from pptx import Presentation
    with patch("agents.tools.powerpoint_output.insert_agent_output_sync"):
        tool = PowerPointOutputTool(slug=slug)
        result = tool._run(
            slides=sample_slides,
            metadata=sample_metadata,
            filename="executive_presentation.pptx",
            agent_name="business_plan_generator",
        )
    prs = Presentation(result)
    all_text = " ".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Case for Change" in all_text
    assert "Our Value Propositions" in all_text


def test_powerpoint_output_appends_pptx_extension(slug, sample_slides, sample_metadata):
    """filename without .pptx extension gets it added automatically."""
    from agents.tools.powerpoint_output import PowerPointOutputTool
    with patch("agents.tools.powerpoint_output.insert_agent_output_sync"):
        tool = PowerPointOutputTool(slug=slug)
        result = tool._run(
            slides=sample_slides,
            metadata=sample_metadata,
            filename="presentation_no_ext",
            agent_name="business_plan_generator",
        )
    assert result.endswith(".pptx")
    assert Path(result).exists()


def test_powerpoint_output_error_on_write_failure(slug, sample_slides, sample_metadata):
    """Returns error string when the file cannot be saved."""
    from agents.tools.powerpoint_output import PowerPointOutputTool
    from pptx.presentation import Presentation as PptxPresentation
    with patch("agents.tools.powerpoint_output.insert_agent_output_sync"), \
         patch.object(PptxPresentation, "save", side_effect=OSError("disk full")):
        tool = PowerPointOutputTool(slug=slug)
        result = tool._run(
            slides=sample_slides,
            metadata=sample_metadata,
            filename="fail.pptx",
            agent_name="business_plan_generator",
        )
    assert result.startswith("Error: render failed")
