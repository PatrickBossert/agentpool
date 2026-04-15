# agents/tools/powerpoint_output.py
from pathlib import Path
from typing import Any
from pydantic import BaseModel, Field
from crewai.tools import BaseTool
from api.config import get_settings
from agents.tools._db import insert_agent_output_sync


class PowerPointOutputToolInput(BaseModel):
    slides: list[dict[str, Any]] = Field(
        description=(
            "List of dicts with 'title' (str), 'content' (str or list[str]), "
            "and 'notes' (str) keys — one per slide."
        )
    )
    metadata: dict[str, Any] = Field(
        description="Dict with keys: org_name, financial_year, sponsor, date."
    )
    filename: str = Field(
        description="Output filename (e.g. 'executive_presentation.pptx'). "
        ".pptx extension added automatically if missing."
    )
    agent_name: str = Field(
        description="Name of the agent producing this output (used for output tracking)."
    )


class PowerPointOutputTool(BaseTool):
    name: str = "PowerPointOutputTool"
    description: str = (
        "Write a presentation to a .pptx file in the project outputs directory. "
        "Pass slides as a list of {title, content, notes} dicts, metadata with "
        "org_name/financial_year/sponsor/date, and a filename. "
        "Returns the absolute file path."
    )
    args_schema: type[BaseModel] = PowerPointOutputToolInput
    slug: str

    def _run(
        self,
        slides: list[dict[str, Any]],
        metadata: dict[str, Any],
        filename: str,
        agent_name: str,
    ) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.dml.color import RGBColor
        except ImportError:
            return "Error: python-pptx not installed — run: pip install python-pptx"

        settings = get_settings()
        outputs_dir = Path(settings.projects_dir) / self.slug / "outputs"
        outputs_dir.mkdir(parents=True, exist_ok=True)

        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"
        file_path = outputs_dir / filename

        _navy = RGBColor(0x1F, 0x39, 0x7D)

        try:
            prs = Presentation()
            # 16:9 widescreen
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # ── Slide 1: Title ─────────────────────────────────────────────
            title_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(title_layout)
            title_shape = slide.shapes.title
            title_shape.text = metadata.get("org_name", "")
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = _navy
                    run.font.name = "Arial"
                    run.font.bold = True
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = (
                    f"Digital Modernisation Business Plan\n"
                    f"{metadata.get('financial_year', '')}  |  "
                    f"{metadata.get('sponsor', '')}  |  "
                    f"{metadata.get('date', '')}"
                )

            # ── Slides 2–N: Content ────────────────────────────────────────
            content_layout = prs.slide_layouts[1]  # Title and Content layout
            for slide_data in slides:
                s = prs.slides.add_slide(content_layout)
                s.shapes.title.text = slide_data.get("title", "")

                # Style title navy
                for para in s.shapes.title.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = _navy
                        run.font.bold = True
                        run.font.name = "Arial"

                # Body content
                body = s.placeholders[1]
                tf = body.text_frame
                tf.clear()
                content = slide_data.get("content", "")
                if isinstance(content, list):
                    for i, bullet in enumerate(content):
                        if i == 0:
                            tf.paragraphs[0].text = str(bullet)
                        else:
                            tf.add_paragraph().text = str(bullet)
                else:
                    tf.paragraphs[0].text = str(content)

                # Speaker notes
                notes_text = slide_data.get("notes", "")
                if notes_text:
                    notes_slide = s.notes_slide
                    notes_slide.notes_text_frame.text = notes_text

            prs.save(file_path)
            insert_agent_output_sync(
                slug=self.slug,
                agent_name=agent_name,
                output_type="pptx",
                file_path=str(file_path),
            )
        except (OSError, ValueError) as e:
            return f"Error: render failed — {e}"

        return str(file_path)
